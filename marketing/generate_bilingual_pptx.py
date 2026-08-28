from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re, os

base_dir = "/opt/data/projects/nephroassist/marketing"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def hex_to_rgb(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

ACCENT = hex_to_rgb("1a3a5c")
BG = hex_to_rgb("f7fafc")
TEXT = hex_to_rgb("1a202c")
LIGHT = hex_to_rgb("e2e8f0")

def add_textbox(slide, left, top, width, height, text, font_size=Pt(18), bold=False, color=TEXT, align=PP_ALIGN.LEFT, font_name="Helvetica Neue"):
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    font = p.font
    font.size = font_size
    font.bold = bold
    font.color.rgb = RGBColor(*color)
    font.name = font_name
    return tf

def add_slide_with_notes(prs, title, body_lines, notes_text, is_cover=False, accent_color=ACCENT):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text
    if is_cover:
        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*accent_color)
        bg.line.fill.background()
        add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12), Inches(1.6), title, font_size=Pt(46), bold=True, color=(255,255,255), align=PP_ALIGN.LEFT)
        add_textbox(slide, Inches(0.5), Inches(3.8), Inches(12), Inches(1.5), "\n".join(body_lines), font_size=Pt(20), bold=False, color=(220,220,220), align=PP_ALIGN.LEFT)
        add_textbox(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.5), "August 2026", font_size=Pt(14), bold=False, color=(180,180,180), align=PP_ALIGN.LEFT)
    else:
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*accent_color)
        bar.line.fill.background()
        add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12), Inches(0.8), title, font_size=Pt(32), bold=True, color=accent_color, align=PP_ALIGN.LEFT)
        body_text = "\n".join(body_lines).strip()
        if not body_text:
            body_text = " "
        add_textbox(slide, Inches(0.5), Inches(1.25), Inches(12), Inches(5.6), body_text, font_size=Pt(16), bold=False, color=TEXT, align=PP_ALIGN.LEFT)
    return slide

def parse_deck(md_path):
    with open(md_path, "r", encoding="utf-8") as f:
        text = f.read()
    slides = []
    current_title = None
    current_body = []
    current_notes = []
    for line in text.splitlines():
        if line.startswith("## Slide") and ":" in line:
            if current_title is not None:
                slides.append((current_title, current_body, current_notes))
            current_title = line.split(":", 1)[1].strip() if ":" in line else line[3:].strip()
            current_body = []
            current_notes = []
        elif line.startswith("# ") and not line.startswith("## "):
            # Document title / cover
            if current_title is not None:
                slides.append((current_title, current_body, current_notes))
            current_title = line[1:].strip()
            current_body = []
            current_notes = []
        elif line.startswith("---"):
            continue
        elif line.startswith("**Sprechernotizen:**") or line.startswith("**Note per il relatore:**"):
            # Start notes
            if current_body:
                pass
        elif line.startswith("**") and "Sprechernotizen" in line:
            pass
        elif line.startswith("**") and "Note per il relatore" in line:
            pass
        else:
            # Check if we are in notes section after body
            # Heuristic: if previous line was notes header, collect notes
            pass
    if current_title is not None:
        slides.append((current_title, current_body, current_notes))
    return slides

# Better parser: split by Slide sections, then detect notes block
def parse_deck_v2(md_path):
    with open(md_path, "r", encoding="utf-8") as f:
        text = f.read()
    parts = re.split(r'\n## Slide \d+:.+\n', text)
    headers = re.findall(r'\n## Slide \d+:\s*(.+?)\n', text)
    slides = []
    # First part before first slide is cover/front matter
    cover_text = parts[0] if parts else ""
    cover_lines = [l.strip() for l in cover_text.splitlines() if l.strip() and not l.startswith("---")]
    cover_title = cover_lines[0] if cover_lines else "NephroAssist"
    cover_sub = "\n".join(cover_lines[1:]) if len(cover_lines) > 1 else ""
    slides.append(("COVER", cover_title, cover_sub.splitlines(), ""))
    for i, header in enumerate(headers):
        body = parts[i+1] if i+1 < len(parts) else ""
        # Split body and notes by "Sprechernotizen" or "Note per il relatore"
        note_split = re.split(r'\n\*\*(?:Sprechernotizen|Note per il relatore):\*\*\n', body)
        body_text = note_split[0] if note_split else body
        notes_text = note_split[1] if len(note_split) > 1 else ""
        # Clean body
        body_lines = []
        in_table = False
        in_code = False
        for ln in body_text.splitlines():
            s = ln.strip()
            if s.startswith("|"):
                continue
            if s.startswith("```"):
                in_code = not in_code
                continue
            if in_code:
                body_lines.append(s)
                continue
            if s.startswith("**") and s.endswith("**") and s.count("**") == 2:
                body_lines.append(s.replace("**", ""))
            elif s.startswith("> "):
                body_lines.append("• " + s[2:])
            elif s.startswith("- ") or s.startswith("* "):
                body_lines.append("• " + s[2:])
            elif re.match(r'^\d+\.\s', s):
                body_lines.append("• " + re.sub(r'^\d+\.\s', '', s))
            elif s == "":
                body_lines.append("")
            else:
                body_lines.append(s.replace("**", ""))
        slides.append(("SLIDE", header, body_lines, notes_text.strip()))
    return slides

def build_pptx(md_path, pptx_path, accent_color=ACCENT):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    parsed = parse_deck_v2(md_path)
    for entry in parsed:
        kind = entry[0]
        if kind == "COVER":
            title = entry[1]
            sub_lines = entry[2]
            # Clean sub
            cleaned_sub = []
            for s in sub_lines:
                if s.startswith("*") and s.endswith("*"):
                    cleaned_sub.append(s.strip("*"))
                else:
                    cleaned_sub.append(s)
            add_slide_with_notes(prs, title, cleaned_sub, "", is_cover=True, accent_color=accent_color)
        else:
            header = entry[1]
            body_lines = entry[2]
            notes = entry[3]
            add_slide_with_notes(prs, header, body_lines, notes, is_cover=False, accent_color=accent_color)
    prs.save(pptx_path)
    print(f"Generated PowerPoint: {pptx_path}")

# German
de_md = os.path.join(base_dir, "nephroassist-pitch-deck-de-v1.md")
de_pptx = os.path.join(base_dir, "nephroassist-pitch-deck-de-v1.pptx")
build_pptx(de_md, de_pptx, accent_color=ACCENT)

# Italian
it_md = os.path.join(base_dir, "nephroassist-pitch-deck-it-v1.md")
it_pptx = os.path.join(base_dir, "nephroassist-pitch-deck-it-v1.pptx")
build_pptx(it_md, it_pptx, accent_color=hex_to_rgb("1a4a3c"))

print("All PPTX files generated.")
