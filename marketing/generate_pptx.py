from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import re, os

base_dir = "/opt/data/projects/nephroassist/marketing"
md_path = os.path.join(base_dir, "nephroassist-pitch-deck.md")
pptx_path = os.path.join(base_dir, "nephroassist-pitch-deck.pptx")

# Read markdown
with open(md_path, "r", encoding="utf-8") as f:
    text = f.read()

# Split into slides by "## Slide N" or "## " headings that look like slide titles
# Our deck uses "## Slide X:" headings. Let's parse them.
# If not present, split by "## " level-2 headings.
slides = []
current_title = None
current_body = []

for line in text.splitlines():
    if line.startswith("## "):
        if current_title is not None:
            slides.append((current_title, "\n".join(current_body)))
        current_title = line[3:].strip()
        current_body = []
    elif line.startswith("# ") and not line.startswith("## "):
        # Title slide
        if current_title is not None:
            slides.append((current_title, "\n".join(current_body)))
        current_title = line[1:].strip()
        current_body = []
    else:
        current_body.append(line)

if current_title is not None:
    slides.append((current_title, "\n".join(current_body)))

# Clean up title slide: the first heading is the document title, body is tagline.
# Let's restructure: first entry is the cover slide.
if slides:
    cover_title = slides[0][0]
    cover_sub = slides[0][1].strip()
    slides = slides[1:]
else:
    cover_title = "NephroAssist"
    cover_sub = ""

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def hex_to_rgb(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

ACCENT = hex_to_rgb("1a3a5c")
ACCENT2 = hex_to_rgb("2c5282")
BG = hex_to_rgb("f7fafc")
TEXT = hex_to_rgb("1a202c")

def add_textbox(slide, left, top, width, height, text, font_size=Pt(18), bold=False, color=TEXT, align=PP_ALIGN.LEFT, font_name="Helvetica Neue"):
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    font = p.font
    font.size = font_size
    font.bold = bold
    font.color.rgb = RGBColor(*color)
    font.name = font_name
    return tf

# Cover slide
blank_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(blank_layout)
# background shape
bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)  # rectangle
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(*ACCENT)
bg.line.fill.background()
add_textbox(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(1.2), cover_title, font_size=Pt(44), bold=True, color=(255,255,255), align=PP_ALIGN.LEFT)
add_textbox(slide, Inches(0.5), Inches(3.6), Inches(12), Inches(1), cover_sub.replace("*", "").replace("Prepared for", "Prepared for"), font_size=Pt(20), bold=False, color=(220,220,220), align=PP_ALIGN.LEFT)
add_textbox(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.5), "August 2026", font_size=Pt(14), bold=False, color=(180,180,180), align=PP_ALIGN.LEFT)

for title, body in slides:
    slide = prs.slides.add_slide(blank_layout)
    # Accent bar at top
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*ACCENT)
    bar.line.fill.background()
    # Title
    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.8), title, font_size=Pt(32), bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
    # Body processing: strip markdown tables/blocks, keep plain text bullet points
    lines = body.splitlines()
    cleaned = []
    in_table = False
    in_code = False
    for ln in lines:
        s = ln.strip()
        if s.startswith("|"):
            continue  # skip tables for pptx (too wide)
        if s.startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            cleaned.append(s)
            continue
        if s.startswith("**") and s.endswith("**"):
            cleaned.append(s.replace("**", ""))
        elif s.startswith("> "):
            cleaned.append("• " + s[2:])
        elif s.startswith("- ") or s.startswith("* "):
            cleaned.append("• " + s[2:])
        elif s.startswith("1. ") or s.startswith("2. ") or s.startswith("3. "):
            cleaned.append("• " + s[3:])
        elif s == "":
            cleaned.append("")
        else:
            cleaned.append(s.replace("**", ""))
    body_text = "\n".join(cleaned).strip()
    if not body_text:
        body_text = " "
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5), body_text, font_size=Pt(16), bold=False, color=TEXT, align=PP_ALIGN.LEFT)

prs.save(pptx_path)
print(f"Generated PowerPoint: {pptx_path}")
